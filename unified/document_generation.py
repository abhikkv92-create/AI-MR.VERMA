#!/usr/bin/env python3
"""
MR.VERMA Document & Code Generation System
===========================================

Complete document generation and code schema management:
- DOCX, PPTX, XLSX document creation
- SQLite database and schema management
- Vector embeddings and search
- Code generation from templates
- Schema generation and validation

Integrated with MR.VERMA Autonomous system.
"""

import os
import sys
import json
import sqlite3
from pathlib import Path
from typing import Dict, List, Optional, Any, Union
from dataclasses import dataclass, field
from datetime import datetime
import subprocess

# Add parent to path
sys.path.insert(0, str(Path(__file__).parent.parent))

# Document generation imports
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    from pptx.enum.text import PP_ALIGN

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    import openpyxl
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment

    XLSX_AVAILABLE = True
except ImportError:
    XLSX_AVAILABLE = False

# Vector embeddings
try:
    import numpy as np
    from unified.vector_store import VectorStore
    VECTOR_AVAILABLE = True
except ImportError:
    VECTOR_AVAILABLE = False

# Template Engine
try:
    from unified.template_engine import TemplateEngine, TemplateEngineType
    TEMPLATE_AVAILABLE = True
except ImportError:
    TEMPLATE_AVAILABLE = False

# Rich UI
try:
    from rich.console import Console
    from rich.panel import Panel

    RICH_AVAILABLE = True
except ImportError:
    RICH_AVAILABLE = False

    class SimpleConsole:
        def print(self, *args, **kwargs):
            print(" ".join(str(a) for a in args))

    Console = SimpleConsole

console = Console()


# ═══════════════════════════════════════════════════════════════════════════════
# DOCUMENT GENERATION SYSTEM
# ═══════════════════════════════════════════════════════════════════════════════


@dataclass
class DocumentConfig:
    """Configuration for document generation"""

    title: str = ""
    author: str = "MR.VERMA"
    subject: str = ""
    output_dir: str = "output/documents"


class DocumentGenerator:
    """Generate DOCX, PPTX, XLSX documents"""

    def __init__(self, config: DocumentConfig = None):
        self.config = config or DocumentConfig()
        os.makedirs(self.config.output_dir, exist_ok=True)

    def create_docx(self, content: Dict, filename: str = None) -> str:
        """Create Word document from content structure"""
        if not DOCX_AVAILABLE:
            return "❌ python-docx not installed. Run: pip install python-docx"

        if not filename:
            filename = f"document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"

        filepath = os.path.join(self.config.output_dir, filename)

        try:
            doc = Document()

            # Add title
            if self.config.title:
                title = doc.add_heading(self.config.title, 0)
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Process content sections
            for section in content.get("sections", []):
                # Add heading
                if "heading" in section:
                    level = section.get("level", 1)
                    doc.add_heading(section["heading"], level=level)

                # Add paragraphs
                if "paragraphs" in section:
                    for para in section["paragraphs"]:
                        p = doc.add_paragraph(para.get("text", ""))
                        if para.get("bold"):
                            p.runs[0].bold = True
                        if para.get("italic"):
                            p.runs[0].italic = True

                # Add bullet points
                if "bullets" in section:
                    for bullet in section["bullets"]:
                        doc.add_paragraph(bullet, style="List Bullet")

                # Add table
                if "table" in section:
                    table_data = section["table"]
                    table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
                    for i, row_data in enumerate(table_data):
                        row = table.rows[i]
                        for j, cell_text in enumerate(row_data):
                            row.cells[j].text = str(cell_text)

            # Save document
            doc.save(filepath)
            return f"✅ Document created: {filepath}"

        except Exception as e:
            return f"❌ Error creating document: {str(e)}"

    def create_pptx(self, slides: List[Dict], filename: str = None) -> str:
        """Create PowerPoint presentation from slides list
        
        Args:
            slides: List of slide dictionaries. Each dict should have:
                - type: 'title', 'content', 'section'
                - title: str
                - content: str (optional)
                - bullets: List[str] (optional)
            filename: Output filename
        """
        if not PPTX_AVAILABLE:
            return "❌ python-pptx not installed. Run: pip install python-pptx"

        if not filename:
            filename = f"presentation_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        filepath = os.path.join(self.config.output_dir, filename)

        try:
            prs = Presentation()

            for slide_data in slides:
                # Choose layout based on type
                slide_type = slide_data.get("type", "content")
                if slide_type == "title":
                    slide_layout = prs.slide_layouts[0]  # Title Slide
                elif slide_type == "section":
                    slide_layout = prs.slide_layouts[2]  # Section Header
                else:
                    slide_layout = prs.slide_layouts[1]  # Title and Content

                slide = prs.slides.add_slide(slide_layout)

                # Set title
                if slide.shapes.title and "title" in slide_data:
                    slide.shapes.title.text = slide_data["title"]

                # Add content (text body)
                if len(slide.placeholders) > 1:
                    body_shape = slide.placeholders[1]
                    tf = body_shape.text_frame
                    
                    # Main content text
                    if "content" in slide_data:
                        tf.text = slide_data["content"]
                    
                    # Bullets
                    if "bullets" in slide_data:
                        # If content was set, append bullets; otherwise start with bullets
                        # If tf.text was set, adding a paragraph makes a new line
                        for bullet in slide_data["bullets"]:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0

            # Save presentation
            prs.save(filepath)
            return f"✅ Presentation created: {filepath}"

        except Exception as e:
            return f"❌ Error creating presentation: {str(e)}"

    def create_xlsx(self, data: Dict, filename: str = None) -> str:
        """Create Excel spreadsheet from data"""
        if not XLSX_AVAILABLE:
            return "❌ openpyxl not installed. Run: pip install openpyxl"

        if not filename:
            filename = f"spreadsheet_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"

        filepath = os.path.join(self.config.output_dir, filename)

        try:
            wb = Workbook()
            ws = wb.active

            # Set sheet title
            if "sheet_name" in data:
                ws.title = data["sheet_name"]

            # Add headers
            if "headers" in data:
                for col, header in enumerate(data["headers"], 1):
                    cell = ws.cell(row=1, column=col, value=header)
                    cell.font = Font(bold=True)
                    cell.fill = PatternFill(
                        start_color="CCCCCC", end_color="CCCCCC", fill_type="solid"
                    )

            # Add data rows
            if "rows" in data:
                for row_idx, row_data in enumerate(data["rows"], 2):
                    for col_idx, value in enumerate(row_data, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)

            # Add formulas
            if "formulas" in data:
                for formula_info in data["formulas"]:
                    row = formula_info.get("row", 1)
                    col = formula_info.get("col", 1)
                    formula = formula_info.get("formula", "")
                    ws.cell(row=row, column=col, value=formula)

            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = column[0].column_letter
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width

            # Save workbook
            wb.save(filepath)
            return f"✅ Spreadsheet created: {filepath}"

        except Exception as e:
            return f"❌ Error creating spreadsheet: {str(e)}"


# ═══════════════════════════════════════════════════════════════════════════════
# DATABASE & SCHEMA MANAGEMENT
# ═══════════════════════════════════════════════════════════════════════════════


class DatabaseManager:
    """Manage SQLite databases and schemas"""

    def __init__(self, db_path: str = "data/mrverma.db"):
        self.db_path = db_path
        os.makedirs(os.path.dirname(db_path), exist_ok=True)
        self.conn = None
        self.cursor = None

    def connect(self):
        """Connect to database"""
        try:
            self.conn = sqlite3.connect(self.db_path)
            self.cursor = self.conn.cursor()
            return True
        except Exception as e:
            return f"❌ Connection error: {str(e)}"

    def create_table(self, table_name: str, schema: Dict) -> str:
        """Create table from schema definition"""
        if not self.conn:
            result = self.connect()
            if isinstance(result, str) and result.startswith("❌"):
                return result

        try:
            # Build CREATE TABLE statement
            columns = []
            for col_name, col_def in schema.items():
                col_type = col_def.get("type", "TEXT")
                constraints = []

                if col_def.get("primary_key"):
                    constraints.append("PRIMARY KEY")
                if col_def.get("autoincrement"):
                    constraints.append("AUTOINCREMENT")
                if col_def.get("not_null"):
                    constraints.append("NOT NULL")
                if col_def.get("unique"):
                    constraints.append("UNIQUE")
                if "default" in col_def:
                    constraints.append(f"DEFAULT {col_def['default']}")

                col_str = f"{col_name} {col_type} {' '.join(constraints)}".strip()
                columns.append(col_str)

            create_sql = (
                f"CREATE TABLE IF NOT EXISTS {table_name} ({', '.join(columns)})"
            )
            self.cursor.execute(create_sql)
            self.conn.commit()

            return f"✅ Table '{table_name}' created successfully"

        except Exception as e:
            return f"❌ Error creating table: {str(e)}"

    def insert_data(self, table_name: str, data: Dict) -> str:
        """Insert data into table"""
        if not self.conn:
            return "❌ Not connected to database"

        try:
            columns = ", ".join(data.keys())
            placeholders = ", ".join(["?" for _ in data])
            values = tuple(data.values())

            insert_sql = f"INSERT INTO {table_name} ({columns}) VALUES ({placeholders})"
            self.cursor.execute(insert_sql, values)
            self.conn.commit()

            return f"✅ Data inserted into '{table_name}'"

        except Exception as e:
            return f"❌ Error inserting data: {str(e)}"

    def query(self, sql: str, params: tuple = ()) -> List:
        """Execute query and return results"""
        if not self.conn:
            return ["❌ Not connected to database"]

        try:
            self.cursor.execute(sql, params)
            return self.cursor.fetchall()
        except Exception as e:
            return [f"❌ Query error: {str(e)}"]

    def get_schema(self, table_name: str = None) -> Dict:
        """Get database schema"""
        if not self.conn:
            return {"error": "Not connected"}

        try:
            if table_name:
                # Get specific table schema
                self.cursor.execute(f"PRAGMA table_info({table_name})")
                columns = self.cursor.fetchall()
                return {
                    table_name: [
                        {
                            "name": col[1],
                            "type": col[2],
                            "not_null": col[3],
                            "default": col[4],
                            "primary_key": col[5],
                        }
                        for col in columns
                    ]
                }
            else:
                # Get all tables
                self.cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
                tables = self.cursor.fetchall()
                schema = {}
                for table in tables:
                    table_name = table[0]
                    self.cursor.execute(f"PRAGMA table_info({table_name})")
                    columns = self.cursor.fetchall()
                    schema[table_name] = [
                        {
                            "name": col[1],
                            "type": col[2],
                            "not_null": col[3],
                            "default": col[4],
                            "primary_key": col[5],
                        }
                        for col in columns
                    ]
                return schema

        except Exception as e:
            return {"error": str(e)}

    def close(self):
        """Close database connection"""
        if self.conn:
            self.conn.close()
            self.conn = None
            self.cursor = None


# ═══════════════════════════════════════════════════════════════════════════════
# VECTOR EMBEDDINGS & SEARCH
# ═══════════════════════════════════════════════════════════════════════════════





# ═══════════════════════════════════════════════════════════════════════════════
# CODE GENERATION FROM TEMPLATES
# ═══════════════════════════════════════════════════════════════════════════════





# ═══════════════════════════════════════════════════════════════════════════════
# UNIFIED DOCUMENT AGENT
# ═══════════════════════════════════════════════════════════════════════════════


class DocumentAgent:
    """Unified agent for document, database, and code generation"""

    def __init__(self):
        self.doc_generator = DocumentGenerator()
        self.db_manager = DatabaseManager()
        
        # Initialize Unified Vector Store
        if VECTOR_AVAILABLE:
            model_name = os.getenv("VECTOR_MODEL", "all-MiniLM-L6-v2")
            backend = os.getenv("VECTOR_BACKEND", "memory")
            try:
                self.vector_store = VectorStore(model_name=model_name, backend_type=backend)
            except Exception as e:
                console.print(f"[yellow]⚠️ Vector Store init failed: {e}. Using in-memory fallback.[/yellow]")
                self.vector_store = VectorStore(model_name="all-MiniLM-L6-v2", backend_type="memory")
        else:
            self.vector_store = None

        # Initialize Unified Template Engine
        if TEMPLATE_AVAILABLE:
            self.template_engine = TemplateEngine()
        else:
            self.template_engine = None

    def create_document(
        self, doc_type: str, content: Dict, filename: str = None
    ) -> str:
        """Create document of specified type and index it"""
        result = ""
        
        # 1. Generate Document
        if doc_type == "docx":
            result = self.doc_generator.create_docx(content, filename)
        elif doc_type == "pptx":
            result = self.doc_generator.create_pptx(content, filename)
        elif doc_type == "xlsx":
            result = self.doc_generator.create_xlsx(content, filename)
        else:
            return f"❌ Unknown document type: {doc_type}"

        # 2. Index in Vector Store (if successful)
        if "✅" in result and self.vector_store:
            try:
                # Extract text for embedding
                text_content = self._extract_text(content)
                if text_content and filename:
                    # Use filename as ID
                    doc_id = filename
                    self.vector_store.add_document(
                        doc_id=doc_id, 
                        text=text_content, 
                        metadata={
                            "type": doc_type, 
                            "path": result.split(": ")[-1].strip(),
                            "created_at": datetime.now().isoformat()
                        }
                    )
                    result += f" + 🧠 Indexed"
            except Exception as e:
                result += f" (⚠️ Indexing failed: {str(e)})"
                
        return result

    def _extract_text(self, content: Dict) -> str:
        """Extract plain text from document content structure"""
        text_parts = []
        
        # Handle sections/paragraphs (common in DOCX)
        if "sections" in content:
            for section in content["sections"]:
                if "heading" in section:
                    text_parts.append(section["heading"])
                if "paragraphs" in section:
                    for p in section["paragraphs"]:
                        text_parts.append(p.get("text", ""))
                if "bullets" in section:
                    text_parts.extend(section["bullets"])
        
        # Handle slides (PPTX) - typically passed as a list of dicts
        if isinstance(content, list):
            for item in content:
                if isinstance(item, dict):
                    if "title" in item:
                        text_parts.append(str(item["title"]))
                    if "content" in item:
                        text_parts.append(str(item["content"]))
                    if "bullets" in item and isinstance(item["bullets"], list):
                        text_parts.extend([str(b) for b in item["bullets"]])
        
        # Handle Excel data
        if "rows" in content:
             for row in content["rows"]:
                 text_parts.extend([str(c) for c in row if c])

        return " ".join(text_parts)

    def create_database_schema(self, schema_definition: Dict) -> str:
        """Create database schema"""
        results = []

        for table_name, schema in schema_definition.items():
            result = self.db_manager.create_table(table_name, schema)
            results.append(result)

        self.db_manager.close()
        return "\n".join(results)

    def generate_code(self, template: str, data: Dict, output: str = None) -> str:
        """Generate code from template using Unified Ecosystem"""
        if not self.template_engine:
            return "❌ Template Engine not available"

        try:
            # Check for built-in templates first
            builtin = self.template_engine._built_in_templates.get(template)
            
            if builtin:
                # Use the built-in template content
                result = self.template_engine.render_template(
                    builtin.template, 
                    data, 
                    builtin.engine
                )
            else:
                # Try to load as file or string?
                # For compatibility, assume 'template' might be a file path if it exists
                if os.path.exists(template):
                     result = self.template_engine.render_file(template, data)
                else:
                    return f"❌ Template '{template}' not found in built-ins or files"

            if output:
                with open(output, "w") as f:
                    f.write(result.content)
                return f"✅ Code generated: {output}"
            
            return result.content
            
        except Exception as e:
            return f"❌ Code generation failed: {str(e)}"

    def get_capabilities(self) -> Dict:
        """Get agent capabilities"""
        caps = {
            "document_types": ["docx", "pptx", "xlsx"],
            "database": "SQLite with schema management",
            "vectors": "Not Available",
            "code_generation": "Unified Template Engine",
            "libraries": {
                "docx": DOCX_AVAILABLE,
                "pptx": PPTX_AVAILABLE,
                "xlsx": XLSX_AVAILABLE,
            },
        }
        
        if self.vector_store:
            caps["vectors"] = f"{self.vector_store.model_name} ({self.vector_store.backend_type})"
            
        return caps


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN INTERFACE
# ═══════════════════════════════════════════════════════════════════════════════


def print_document_banner():
    """Print document generation banner"""
    banner = """
===============================================================
           MR.VERMA Document & Code Generation
===============================================================
  DOCX  |  PPTX  |  XLSX  |  SQLite  |  Vectors
===============================================================
    """
    console.print(banner)


def demo_document_generation():
    """Demonstrate document generation capabilities"""
    print_document_banner()

    agent = DocumentAgent()

    # Show capabilities
    console.print("\n[bold cyan]Agent Capabilities:[/bold cyan]")
    caps = agent.get_capabilities()
    for key, value in caps.items():
        console.print(f"  {key}: {value}")

    # Demo 1: Create DOCX
    console.print("\n[bold green]Demo 1: Creating Word Document[/bold green]")
    doc_content = {
        "sections": [
            {
                "heading": "Introduction",
                "level": 1,
                "paragraphs": [
                    {"text": "This is an automatically generated document."},
                    {"text": "Created by MR.VERMA Document Agent.", "bold": True},
                ],
                "bullets": [
                    "Feature 1: Automatic generation",
                    "Feature 2: Multiple formats",
                    "Feature 3: Easy to use",
                ],
            }
        ]
    }
    result = agent.create_document("docx", doc_content, "demo_document.docx")
    console.print(f"  {result}")

    # Demo 2: Create XLSX
    console.print("\n[bold green]Demo 2: Creating Excel Spreadsheet[/bold green]")
    xlsx_data = {
        "sheet_name": "Project Data",
        "headers": ["Task", "Status", "Priority", "Assignee"],
        "rows": [
            ["Design UI", "Complete", "High", "Alice"],
            ["Backend API", "In Progress", "High", "Bob"],
            ["Testing", "Pending", "Medium", "Charlie"],
            ["Documentation", "Pending", "Low", "Diana"],
        ],
    }
    result = agent.create_document("xlsx", xlsx_data, "demo_spreadsheet.xlsx")
    console.print(f"  {result}")

    # Demo 3: Create Database
    console.print("\n[bold green]Demo 3: Creating Database Schema[/bold green]")
    schema = {
        "users": {
            "id": {"type": "INTEGER", "primary_key": True, "autoincrement": True},
            "name": {"type": "TEXT", "not_null": True},
            "email": {"type": "TEXT", "unique": True},
            "created_at": {"type": "TEXT", "default": "CURRENT_TIMESTAMP"},
        },
        "projects": {
            "id": {"type": "INTEGER", "primary_key": True, "autoincrement": True},
            "name": {"type": "TEXT", "not_null": True},
            "user_id": {"type": "INTEGER", "not_null": True},
            "status": {"type": "TEXT", "default": '"active"'},
        },
    }
    result = agent.create_database_schema(schema)
    console.print(f"  {result}")

    console.print("\n[bold cyan]Document Generation Demo Complete![/bold cyan]")


if __name__ == "__main__":
    demo_document_generation()
